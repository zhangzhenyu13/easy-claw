"""
文件处理器Skill实现
内联了FileProcessor的核心逻辑和format_file_context函数
"""
import os
import tempfile
import base64
import mimetypes
from pathlib import Path
from typing import List, Optional, Tuple, Union
from io import BytesIO
import json
import logging

from tir_agent.skills.base import BaseSkill, SkillMetadata

# 配置日志
logger = logging.getLogger(__name__)


def format_file_context(processed_files: List[dict]) -> str:
    """
    将处理后的文件信息格式化为上下文字符串 - 内联自原file_processor

    Args:
        processed_files: 处理后的文件列表

    Returns:
        str: 格式化后的上下文字符串
    """
    context_parts = []

    for pf in processed_files:
        header = f"\n--- 文件: {pf['file_name']} (类型: {pf['file_type']}) ---\n"

        if pf.get('error'):
            context_parts.append(f"{header}[处理错误: {pf['error']}]")
            continue

        content = pf.get('content', '')
        if content:
            context_parts.append(f"{header}{content}")

        if pf.get('images'):
            img_count = len(pf['images'])
            context_parts.append(f"{header}[包含 {img_count} 张图片]")

    return "\n".join(context_parts)


class FileProcessorSkill(BaseSkill):
    """文件处理器Skill - 直接内联文件处理逻辑"""

    # 支持的图片格式
    IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff'}
    # 支持的文档格式
    DOCUMENT_EXTENSIONS = {'.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.txt', '.md'}

    def __init__(self, metadata: SkillMetadata, config: dict = None):
        super().__init__(metadata, config)
        # 从配置获取临时目录
        self.temp_dir = config.get("temp_dir") if config else tempfile.gettempdir()
        os.makedirs(self.temp_dir, exist_ok=True)

    def is_image(self, file_path: str) -> bool:
        """判断文件是否为图片"""
        ext = Path(file_path).suffix.lower()
        return ext in self.IMAGE_EXTENSIONS

    def is_document(self, file_path: str) -> bool:
        """判断文件是否为文档"""
        ext = Path(file_path).suffix.lower()
        return ext in self.DOCUMENT_EXTENSIONS

    def get_file_type(self, file_path: str) -> str:
        """获取文件类型"""
        ext = Path(file_path).suffix.lower()
        mime_type, _ = mimetypes.guess_type(file_path)

        if ext in self.IMAGE_EXTENSIONS:
            return "image"
        elif ext == '.pdf':
            return "pdf"
        elif ext in {'.docx', '.doc'}:
            return "word"
        elif ext in {'.xlsx', '.xls'}:
            return "excel"
        elif ext in {'.pptx', '.ppt'}:
            return "ppt"
        elif ext in {'.txt', '.md'}:
            return "text"
        else:
            return "unknown"

    def process_file(self, file_path: str) -> dict:
        """
        处理单个文件，返回处理结果 - 内联自原FileProcessor

        Returns:
            dict: 包含文件类型、内容和元数据的字典
        """
        logger.info(f"📄 开始处理文件: {file_path}")
        file_type = self.get_file_type(file_path)
        logger.info(f"   文件类型: {file_type}")

        result = {
            "file_path": file_path,
            "file_name": Path(file_path).name,
            "file_type": file_type,
            "content": None,
            "images": [],
            "metadata": {}
        }

        try:
            if file_type == "image":
                logger.info(f"   处理图片文件...")
                result["images"] = [self._process_image(file_path)]
                result["content"] = f"[图片文件: {result['file_name']}]"
                logger.info(f"   ✅ 图片处理完成，尺寸: {result['images'][0]['width']}x{result['images'][0]['height']}")

            elif file_type == "pdf":
                logger.info(f"   处理PDF文件...")
                text, images = self._process_pdf(file_path)
                result["content"] = text
                result["images"] = images
                result["metadata"]["is_scanned"] = len(text.strip()) < 100 and len(images) > 0
                logger.info(f"   ✅ PDF处理完成 - 文本长度: {len(text)}, 图片数: {len(images)}, 是否扫描件: {result['metadata']['is_scanned']}")

            elif file_type == "word":
                logger.info(f"   处理Word文档...")
                result["content"] = self._process_word(file_path)
                logger.info(f"   ✅ Word处理完成，内容长度: {len(result['content'] or '')}")

            elif file_type == "excel":
                logger.info(f"   处理Excel文件...")
                result["content"] = self._process_excel(file_path)
                logger.info(f"   ✅ Excel处理完成，内容长度: {len(result['content'] or '')}")

            elif file_type == "ppt":
                logger.info(f"   处理PPT文件...")
                result["content"] = self._process_ppt(file_path)
                logger.info(f"   ✅ PPT处理完成，内容长度: {len(result['content'] or '')}")

            elif file_type == "text":
                logger.info(f"   处理文本文件...")
                result["content"] = self._process_text(file_path)
                logger.info(f"   ✅ 文本处理完成，内容长度: {len(result['content'] or '')}")

            else:
                logger.warning(f"   ⚠️ 不支持的文件类型: {file_type}")
                result["content"] = f"[不支持的文件类型: {result['file_name']}]"

        except Exception as e:
            logger.error(f"   ❌ 处理文件失败 {file_path}: {e}", exc_info=True)
            result["error"] = str(e)
            result["content"] = f"[文件处理失败: {str(e)}]"

        return result

    def _process_image(self, file_path: str) -> dict:
        """处理图片文件，返回图片信息"""
        from PIL import Image

        with Image.open(file_path) as img:
            return {
                "path": file_path,
                "width": img.width,
                "height": img.height,
                "format": img.format,
                "mode": img.mode,
                "base64": self._image_to_base64(img)
            }

    def _image_to_base64(self, img) -> str:
        """将图片转换为base64编码"""
        buffer = BytesIO()
        # 转换为RGB模式（如果是RGBA或其他模式）
        if img.mode in ('RGBA', 'P'):
            img = img.convert('RGB')
        img.save(buffer, format='JPEG')
        return base64.b64encode(buffer.getvalue()).decode('utf-8')

    def _process_pdf(self, file_path: str) -> Tuple[str, List[dict]]:
        """
        处理PDF文件
        对于扫描件PDF，转换为图片
        对于电子文档PDF，提取文本

        Returns:
            Tuple[str, List[dict]]: (文本内容, 图片列表)
        """
        text_content = ""
        images = []

        # 尝试提取文本
        try:
            logger.info(f"      尝试提取PDF文本...")
            text_content = self._extract_pdf_text(file_path)
            logger.info(f"      提取到文本长度: {len(text_content)}")
        except Exception as e:
            logger.warning(f"      PDF文本提取失败: {e}")

        # 如果文本很少，可能是扫描件，转换为图片
        if len(text_content.strip()) < 100:
            logger.info(f"      文本较少，可能是扫描件，尝试转换为图片...")
            try:
                images = self._pdf_to_images(file_path)
                logger.info(f"      PDF转图片完成，共 {len(images)} 页")
            except Exception as e:
                logger.warning(f"      PDF转图片失败: {e}")
        else:
            logger.info(f"      文本充足，跳过图片转换")

        return text_content, images

    def _extract_pdf_text(self, file_path: str) -> str:
        """提取PDF文本内容"""
        try:
            from pypdf import PdfReader

            reader = PdfReader(file_path)
            logger.info(f"      PDF页数: {len(reader.pages)}")
            text_parts = []
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text:
                    text_parts.append(text)
                    logger.debug(f"      第{i+1}页文本长度: {len(text)}")
            return "\n\n".join(text_parts)
        except ImportError:
            logger.warning("      pypdf未安装，无法提取PDF文本")
            return ""
        except Exception as e:
            logger.error(f"      PDF文本提取出错: {e}")
            return ""

    def _pdf_to_images(self, file_path: str) -> List[dict]:
        """将PDF转换为图片列表"""
        try:
            from pdf2image import convert_from_path

            logger.info(f"      正在将PDF转换为图片...")
            images = convert_from_path(file_path)
            logger.info(f"      转换完成，共 {len(images)} 页")
            result = []

            for i, img in enumerate(images):
                # 保存临时图片
                temp_path = os.path.join(
                    self.temp_dir,
                    f"{Path(file_path).stem}_page_{i+1}.jpg"
                )
                img.save(temp_path, 'JPEG')
                logger.info(f"      保存第{i+1}页图片: {temp_path}")

                result.append({
                    "path": temp_path,
                    "page": i + 1,
                    "width": img.width,
                    "height": img.height,
                    "base64": self._image_to_base64(img)
                })

            return result
        except ImportError as e:
            logger.warning(f"      pdf2image未安装，无法将PDF转换为图片: {e}")
            return []
        except Exception as e:
            logger.error(f"      PDF转图片失败: {e}", exc_info=True)
            return []

    def _process_word(self, file_path: str) -> str:
        """处理Word文档"""
        try:
            from docx import Document

            doc = Document(file_path)
            text_parts = []

            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # 提取表格内容
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells]
                    table_text.append(" | ".join(row_text))
                text_parts.append("\n".join(table_text))

            return "\n\n".join(text_parts)
        except ImportError:
            return "[python-docx未安装，无法处理Word文档]"
        except Exception as e:
            return f"[Word文档处理失败: {str(e)}]"

    def _process_excel(self, file_path: str) -> str:
        """处理Excel文件"""
        try:
            import openpyxl

            wb = openpyxl.load_workbook(file_path, data_only=True)
            text_parts = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"=== Sheet: {sheet_name} ===")

                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) if cell is not None else "" for cell in row]
                    if any(cell.strip() for cell in row_text):
                        text_parts.append(" | ".join(row_text))

            return "\n".join(text_parts)
        except ImportError:
            return "[openpyxl未安装，无法处理Excel文件]"
        except Exception as e:
            return f"[Excel文件处理失败: {str(e)}]"

    def _process_ppt(self, file_path: str) -> str:
        """处理PPT文件"""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            text_parts = []

            for i, slide in enumerate(prs.slides, 1):
                text_parts.append(f"=== 幻灯片 {i} ===")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)

            return "\n\n".join(text_parts)
        except ImportError:
            return "[python-pptx未安装，无法处理PPT文件]"
        except Exception as e:
            return f"[PPT文件处理失败: {str(e)}]"

    def _process_text(self, file_path: str) -> str:
        """处理纯文本文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            with open(file_path, 'r', encoding='gbk') as f:
                return f.read()

    def process_files(self, file_paths: List[str]) -> List[dict]:
        """批量处理文件 - agent.py 第349行使用"""
        return [self.process_file(fp) for fp in file_paths]

    def cleanup(self):
        """清理临时文件"""
        import shutil
        try:
            if os.path.exists(self.temp_dir):
                for file in os.listdir(self.temp_dir):
                    file_path = os.path.join(self.temp_dir, file)
                    try:
                        if os.path.isfile(file_path):
                            os.remove(file_path)
                    except Exception as e:
                        logger.warning(f"清理临时文件失败 {file_path}: {e}")
        except Exception as e:
            logger.warning(f"清理临时目录失败: {e}")

    def execute(self, **params) -> str:
        """
        Skill接口方法 - 处理文件

        Args:
            file_path: 文件路径
            output_format: 输出格式（text/markdown/json）

        Returns:
            str: 处理结果
        """
        try:
            # 验证参数
            validated_params = self.validate_params(params)
            file_path = validated_params.get("file_path", "")
            output_format = validated_params.get("output_format", "text")

            # 处理文件
            result = self.process_file(file_path)

            # 根据输出格式返回结果
            if output_format == "json":
                # 返回JSON格式
                return json.dumps(result, ensure_ascii=False, indent=2)

            elif output_format == "markdown":
                # 返回Markdown格式
                return self._format_as_markdown(result)

            else:  # text
                # 返回纯文本格式
                return self._format_as_text(result)

        except ValueError as e:
            return f"参数错误: {str(e)}"
        except Exception as e:
            return f"文件处理失败: {str(e)}"

    def _format_as_text(self, result: dict) -> str:
        """格式化为纯文本"""
        lines = []
        lines.append(f"文件: {result.get('file_name', 'Unknown')}")
        lines.append(f"类型: {result.get('file_type', 'Unknown')}")

        if result.get('error'):
            lines.append(f"错误: {result['error']}")
            return "\n".join(lines)

        content = result.get('content', '')
        if content:
            lines.append(f"\n内容:\n{content}")

        images = result.get('images', [])
        if images:
            lines.append(f"\n[包含 {len(images)} 张图片]")

        metadata = result.get('metadata', {})
        if metadata:
            lines.append(f"\n元数据: {json.dumps(metadata, ensure_ascii=False)}")

        return "\n".join(lines)

    def _format_as_markdown(self, result: dict) -> str:
        """格式化为Markdown"""
        lines = []
        lines.append(f"## 📄 {result.get('file_name', 'Unknown')}")
        lines.append(f"**类型:** `{result.get('file_type', 'Unknown')}`")

        if result.get('error'):
            lines.append(f"\n> ⚠️ **错误:** {result['error']}")
            return "\n".join(lines)

        content = result.get('content', '')
        if content:
            lines.append(f"\n### 内容\n{content}")

        images = result.get('images', [])
        if images:
            lines.append(f"\n### 图片 ({len(images)}张)")
            for img in images:
                if 'page' in img:
                    lines.append(f"- 第{img['page']}页: {img['width']}x{img['height']}")
                else:
                    lines.append(f"- 图片: {img.get('width', 0)}x{img.get('height', 0)}")

        metadata = result.get('metadata', {})
        if metadata:
            lines.append(f"\n### 元数据")
            for key, value in metadata.items():
                lines.append(f"- **{key}:** {value}")

        return "\n".join(lines)
